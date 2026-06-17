"""
@author: { FALCONRY SOLUTIONS }
@description: Pure file generation functions — each takes structured data and returns
              raw bytes ready for streaming to the client or uploading to object store.

Supported formats:
  - CSV / XLS / HTML  : tabular data via Pandas
  - TXT / MD          : plain text (CrewAI agent output used as-is)
  - DOCX (Word)       : rich document via python-docx (docxtpl available for template-based generation)
  - PDF               : generated indirectly via DOCX → docx2pdf (preserves layout)
  - PPTX              : slide deck via python-pptx; LLM produces [{title, content}] slide dicts
"""

import io
import os
import tempfile
from typing import Any, Dict, List


# ── Tabular formats (CSV / XLS / HTML) ───────────────────────────────────────

def generate_csv(records: List[Dict[str, Any]]) -> bytes:
    """Convert a list of flat dicts to a UTF-8 encoded CSV."""
    import pandas as pd

    df = pd.DataFrame(records)
    buffer = io.StringIO()
    df.to_csv(buffer, index=False)
    return buffer.getvalue().encode("utf-8")


def generate_xls(records: List[Dict[str, Any]], sheet_name: str = "Report") -> bytes:
    """Convert a list of flat dicts to an Excel (.xlsx) workbook."""
    import pandas as pd

    df = pd.DataFrame(records)
    buffer = io.BytesIO()
    with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name=sheet_name)
    return buffer.getvalue()


def generate_html(records: List[Dict[str, Any]], title: str = "Report") -> bytes:
    """Convert a list of flat dicts to a styled HTML page containing a data table."""
    import pandas as pd

    df = pd.DataFrame(records)
    table_html = df.to_html(index=False, border=0, classes="data-table")

    page = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <title>{title}</title>
  <style>
    body {{ font-family: Arial, sans-serif; margin: 40px; color: #333; }}
    h1   {{ border-bottom: 2px solid #4a4a8a; padding-bottom: 8px; }}
    .data-table {{ border-collapse: collapse; width: 100%; margin-top: 20px; }}
    .data-table th {{ background: #4a4a8a; color: #fff; padding: 10px 14px; text-align: left; }}
    .data-table td {{ padding: 8px 14px; border-bottom: 1px solid #ddd; }}
    .data-table tr:hover td {{ background: #f5f5ff; }}
  </style>
</head>
<body>
  <h1>{title}</h1>
  {table_html}
</body>
</html>"""

    return page.encode("utf-8")


# ── Plain-text formats (TXT / MD) ─────────────────────────────────────────────
# CrewAI agent text output is used directly — no additional processing needed.

def generate_txt(content: str) -> bytes:
    """Encode a plain-text string as a UTF-8 .txt file."""
    return content.encode("utf-8")


def generate_md(content: str) -> bytes:
    """Encode a Markdown string as a UTF-8 .md file."""
    return content.encode("utf-8")


# ── Word document (DOCX) ──────────────────────────────────────────────────────

def generate_docx(title: str, sections: List[Dict[str, str]]) -> bytes:
    """
    Build a Word document from a title and a list of section dicts.

    Each section dict has two optional keys:
      "heading"  — rendered as a Heading 1
      "content"  — rendered as a normal paragraph

    Uses python-docx directly (the same library that docxtpl wraps) so we can
    build documents programmatically without requiring a pre-existing .docx template.
    For template-based generation (when design teams supply branded .docx templates)
    import and use docxtpl.DocxTemplate instead.
    """
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.oxml.ns import qn
    import docx.oxml

    doc = Document()

    # ── Title ──
    title_para = doc.add_heading(title, level=0)
    title_para.runs[0].font.color.rgb = RGBColor(0x4A, 0x4A, 0x8A)

    # ── Sections ──
    for section in sections:
        heading = section.get("heading", "").strip()
        content = section.get("content", "").strip()

        if heading:
            doc.add_heading(heading, level=1)
        if content:
            doc.add_paragraph(content)

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


# ── PDF (via DOCX → docx2pdf) ────────────────────────────────────────────────

def generate_pdf(title: str, sections: List[Dict[str, str]]) -> bytes:
    """
    Generate a PDF by first building a DOCX in memory, writing it to a temp file,
    converting via docx2pdf (uses MS Word on macOS/Windows, LibreOffice on Linux),
    then reading the result back into memory and cleaning up temp files.
    """
    from docx2pdf import convert

    docx_bytes = generate_docx(title, sections)

    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp_docx:
        tmp_docx.write(docx_bytes)
        tmp_docx_path = tmp_docx.name

    tmp_pdf_path = tmp_docx_path.replace(".docx", ".pdf")

    try:
        convert(tmp_docx_path, tmp_pdf_path)
        with open(tmp_pdf_path, "rb") as f:
            return f.read()
    finally:
        if os.path.exists(tmp_docx_path):
            os.unlink(tmp_docx_path)
        if os.path.exists(tmp_pdf_path):
            os.unlink(tmp_pdf_path)


# ── PowerPoint (PPTX) ────────────────────────────────────────────────────────

def generate_pptx(presentation_title: str, slides: List[Dict[str, str]]) -> bytes:
    """
    Build a PowerPoint presentation from a title and a list of slide dicts.

    Each slide dict must have:
      "title"   — slide title (string)
      "content" — slide body text (string)

    The LLM is expected to return slides as an array of these dicts so that
    each dictionary maps directly to one slide here.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()

    # ── Cover slide ──
    cover_layout = prs.slide_layouts[0]  # Title Slide layout
    cover = prs.slides.add_slide(cover_layout)
    cover.shapes.title.text = presentation_title
    if cover.placeholders[1]:
        cover.placeholders[1].text = "Generated by Falconry AI"

    # ── Content slides ──
    content_layout = prs.slide_layouts[1]  # Title and Content layout
    for slide_data in slides:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = slide_data.get("title", "")

        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        tf.word_wrap = True
        tf.text = slide_data.get("content", "")

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()
